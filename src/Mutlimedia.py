import os
from dotenv import load_dotenv
import gradio as gr
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.chains import RetrievalQA
from langchain.prompts import PromptTemplate
import uuid

# ===========================================
# 🔧 CONFIGURATION
# ===========================================
load_dotenv()

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 100
MAX_TOKENS = 4096
MODEL_NAME = "gpt-4o-mini"
TEMPERATURE = 0.4
CHROMA_BASE_DIR = "chroma_store"

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    raise ValueError("Please set OPENAI_API_KEY in your .env file")

llm = ChatOpenAI(
    api_key=OPENAI_API_KEY,
    model_name=MODEL_NAME,
    temperature=TEMPERATURE,
    max_tokens=MAX_TOKENS
)

PROMPT = PromptTemplate(
    template="""Context: {context}

Question: {question}

Answer the question concisely based on the given context. 
If the context doesn't contain relevant information, say "I don't have enough information to answer that question." """,
    input_variables=["context", "question"]
)

# ===========================================
# 🧩 DOCUMENT EXTRACTORS (TEXT ONLY)
# ===========================================

def process_pdf(pdf_file):
    pdf_reader = PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    splitter = CharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    return splitter.split_text(text)


def process_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    splitter = CharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    return splitter.split_text(text)


def process_docx(docx_file):
    doc = Document(docx_file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    splitter = CharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    return splitter.split_text(text)

# ===========================================
# 🧠 RAG PIPELINE
# ===========================================

def create_vectorstore(texts, collection_name):
    embeddings = OpenAIEmbeddings()
    persist_dir = os.path.join(CHROMA_BASE_DIR, collection_name)
    os.makedirs(persist_dir, exist_ok=True)

    if os.path.exists(os.path.join(persist_dir, "chroma.sqlite3")):
        vectorstore = Chroma(persist_directory=persist_dir, embedding_function=embeddings)
    else:
        vectorstore = Chroma.from_texts(
            texts=texts,
            embedding=embeddings,
            persist_directory=persist_dir
        )
        vectorstore.persist()

    return vectorstore


def process_file_and_query(uploaded_file, query):
    ext = os.path.splitext(uploaded_file)[1].lower()
    file_id = str(uuid.uuid4())

    if ext == ".pdf":
        texts = process_pdf(uploaded_file)
    elif ext == ".pptx":
        texts = process_pptx(uploaded_file)
    elif ext == ".docx":
        texts = process_docx(uploaded_file)
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, PPTX, or DOCX.")

    vectorstore = create_vectorstore(texts, collection_name=file_id)

    qa = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=vectorstore.as_retriever(),
        chain_type_kwargs={"prompt": PROMPT}
    )

    result = qa.run(query)
    return result, len(texts)

# ===========================================
# 🎛 GRADIO INTERFACE
# ===========================================

def gradio_interface(uploaded_file, query):
    result, num_chunks = process_file_and_query(uploaded_file.name, query)
    log = f"File processed successfully. Text chunks: {num_chunks}"
    return result, log


iface = gr.Interface(
    fn=gradio_interface,
    inputs=[
        gr.File(label="Upload PDF, PPTX, or DOCX"),
        gr.Textbox(label="Enter your question")
    ],
    outputs=[
        gr.Textbox(label="Answer"),
        gr.Textbox(label="Processing Log")
    ],
    title="📚 Text-Only RAG System",
    description="Upload a PDF, PPTX, or DOCX and ask questions about its text content."
)

iface.launch()
