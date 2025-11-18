# rag/loaders.py
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import CharacterTextSplitter

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 100

splitter = CharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP
)

def load_pdf(file_path):
    pdf_reader = PdfReader(file_path)
    text = ""
    for page in pdf_reader.pages:
        p = page.extract_text()
        if p:
            text += p + "\n"
    return splitter.split_text(text)

def load_docx(file_path):
    doc = Document(file_path)
    text = "\n".join([p.text for p in doc.paragraphs])
    return splitter.split_text(text)

def load_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return splitter.split_text(text)

def load_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return load_pdf(file_path)
    if ext == ".docx":
        return load_docx(file_path)
    if ext == ".pptx":
        return load_pptx(file_path)

    raise ValueError("Unsupported file type")
