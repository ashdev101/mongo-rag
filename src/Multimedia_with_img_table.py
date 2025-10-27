import os
from dotenv import load_dotenv
import gradio as gr
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import Chroma
from langchain.chains import RetrievalQA, LLMChain
from langchain.prompts import PromptTemplate
import pymupdf as fitz
from PIL import Image
import io
import pandas as pd
import uuid

# ===========================================
# 🔧 CONFIGURATION
# ===========================================
load_dotenv()

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 100
MAX_TOKENS = 4096
MODEL_NAME = "gpt-4o-mini"
TEMPERATURE = 0.4
CHROMA_BASE_DIR = "chroma_store"

OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
if not OPENAI_API_KEY:
    raise ValueError("Please set OPENAI_API_KEY in your .env file")

# Initialize LLM
llm = ChatOpenAI(
    api_key=OPENAI_API_KEY,
    model_name=MODEL_NAME,
    temperature=TEMPERATURE,
    max_tokens=MAX_TOKENS
)

# QA Prompt Template
PROMPT = PromptTemplate(
    template="""Context: {context}

Question: {question}

Answer the question concisely based on the given context. 
If the context doesn't contain relevant information, say "I don't have enough information to answer that question."
If the question is about images or tables, refer to them specifically in your answer.""",
    input_variables=["context", "question"]
)

# ===========================================
# 🧩 DOCUMENT EXTRACTORS
# ===========================================

# --- PDF ---
def process_pdf(pdf_file):
    pdf_reader = PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"

    text_splitter = CharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    texts = text_splitter.split_text(text)
    return texts

def extract_images_and_tables(pdf_file):
    doc = fitz.open(pdf_file)
    images, tables = [], []

    for page_num in range(len(doc)):
        page = doc[page_num]

        # Extract images
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(io.BytesIO(image_bytes))
            images.append((f"Page {page_num + 1}, Image {img_index + 1}", image))

        # Extract tables (PyMuPDF built-in)
        tables_on_page = page.find_tables()
        for table_index, table in enumerate(tables_on_page):
            df = pd.DataFrame(table.extract())
            tables.append((f"Page {page_num + 1}, Table {table_index + 1}", df))

    return images, tables


# --- PPTX ---
def process_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text, images, tables = "", [], []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape_num, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                img = Image.open(io.BytesIO(image_bytes))
                images.append((f"Slide {slide_num}, Image {shape_num}", img))
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                table_data = [[cell.text for cell in row.cells] for row in table.rows]
                df = pd.DataFrame(table_data)
                tables.append((f"Slide {slide_num}, Table {shape_num}", df))

        text += "\n".join(slide_text) + "\n"

    text_splitter = CharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    texts = text_splitter.split_text(text)
    return texts, images, tables


# --- DOCX ---
def process_docx(docx_file):
    doc = Document(docx_file)
    text, images, tables = "", [], []

    # Extract text
    for para in doc.paragraphs:
        text += para.text + "\n"

    # Extract images
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_data = rel.target_part.blob
            img = Image.open(io.BytesIO(image_data))
            images.append(("Embedded Image", img))

    # Extract tables
    for table_num, table in enumerate(doc.tables, start=1):
        table_data = [[cell.text for cell in row.cells] for row in table.rows]
        df = pd.DataFrame(table_data)
        tables.append((f"Table {table_num}", df))

    text_splitter = CharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    texts = text_splitter.split_text(text)
    return texts, images, tables


# ===========================================
# 🧠 RAG PIPELINE (Using Chroma)
# ===========================================
def create_embeddings_and_vectorstore(texts, collection_name):
    """Create or load persistent Chroma vector store."""
    embeddings = OpenAIEmbeddings()

    persist_dir = os.path.join(CHROMA_BASE_DIR, collection_name)
    os.makedirs(persist_dir, exist_ok=True)

    # Reuse if collection already exists
    if os.path.exists(os.path.join(persist_dir, "chroma.sqlite3")):
        vectorstore = Chroma(
            persist_directory=persist_dir,
            embedding_function=embeddings
        )
    else:
        vectorstore = Chroma.from_texts(
            texts=texts,
            embedding=embeddings,
            persist_directory=persist_dir
        )
        vectorstore.persist()

    return vectorstore

def expand_query(query: str, llm: ChatOpenAI) -> str:
    """Expand the user query with semantically related terms."""
    prompt = PromptTemplate(
        input_variables=["query"],
        template="""Given the following query, generate 3–5 related terms or phrases that could be relevant. 
        Separate them with commas.

        Query: {query}

        Related terms:"""
    )
    chain = LLMChain(llm=llm, prompt=prompt)
    response = chain.run(query)
    expanded_terms = [term.strip() for term in response.split(',')]
    return f"{query} {' '.join(expanded_terms)}"

def rag_pipeline(query, qa_chain, vectorstore, images, tables):
    expanded_query = expand_query(query, llm)
    relevant_docs = vectorstore.similarity_search_with_score(expanded_query, k=3)

    context, log = "", f"Query Expansion:\nOriginal: {query}\nExpanded: {expanded_query}\n\nRelevant Chunks:\n"
    for i, (doc, score) in enumerate(relevant_docs, 1):
        context += doc.page_content + "\n\n"
        log += f"Chunk {i} (Score: {score:.4f}): {doc.page_content[:200]}...\n\n"

    context += f"Number of images: {len(images)}\nNumber of tables: {len(tables)}\n"
    response = qa_chain.invoke({"query": query})
    return response['result'], log


# ===========================================
# 🧩 UNIVERSAL FILE HANDLER
# ===========================================
def process_file_and_query(uploaded_file, query):
    ext = os.path.splitext(uploaded_file)[1].lower()
    file_id = str(uuid.uuid4())  # Unique collection for each upload

    if ext == ".pdf":
        texts = process_pdf(uploaded_file)
        images, tables = extract_images_and_tables(uploaded_file)
    elif ext == ".pptx":
        texts, images, tables = process_pptx(uploaded_file)
    elif ext == ".docx":
        texts, images, tables = process_docx(uploaded_file)
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, PPTX, or DOCX.")

    vectorstore = create_embeddings_and_vectorstore(texts, collection_name=file_id)

    qa = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=vectorstore.as_retriever(),
        chain_type_kwargs={"prompt": PROMPT}
    )

    result, log = rag_pipeline(query, qa, vectorstore, images, tables)
    return result, len(texts), len(images), len(tables), log


# ===========================================
# 🎛 GRADIO INTERFACE
# ===========================================
def gradio_interface(uploaded_file, query):
    result, num_chunks, num_images, num_tables, chunks_log = process_file_and_query(uploaded_file.name, query)
    log = (
        f"File processed successfully.\n"
        f"Text chunks: {num_chunks}\nImages: {num_images}\nTables: {num_tables}\n\n{chunks_log}"
    )
    return result, log

# def main():
iface = gr.Interface(
    fn=gradio_interface,
    inputs=[gr.File(label="Upload PDF, PPTX, or DOCX"), gr.Textbox(label="Enter your question")],
    outputs=[gr.Textbox(label="Answer"), gr.Textbox(label="Processing Log")],
    title="📚 Multi-Format RAG System",
    description="Upload a PDF, PPTX, or DOCX and ask questions about its content (including images & tables)."
)
iface.launch()

# if __name__ == "__main__":
#     main()
